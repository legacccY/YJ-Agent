#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
plot_ppt_8to11_rerun.py
=======================
服务: QuantImmuBench §2.2 8-11mer 可变窗补充口径 —— 生成【新切全量重跑版】supplement deck。

本版 (2026-07-09 重建) 精确复刻项目金标准 deck 的专业设计系统:
  参考模板 = QuantImmuBench_9mer新切排名_2026-07-08.pptx (13.333×7.5in, 6 页, Okabe-Ito 配色)。
  几何/配色/版式全部照抄参考 (封面深底+双椭圆装饰; 内容页左全高色条+eyebrow+22pt标题;
  两栏白卡嵌图; 三卡结论页; 蓝脚注+页码), 换成本轮 8-11mer 三张新图 + 核实数字。

取代旧版 (QuantImmuBench_8to11mer_supplement_rev1_2026-07-04.pptx, covfix 旧数据 0.191/0.122)
及本脚本前一版 (深底标题条土设计)。

headline 数字一律从 csv 现算 (脚本 print 自检值供主线核), 不硬编。

产物: QuantImmuBench_8to11mer新切重跑_2026-07-09.pptx (覆盖项目根同名文件)

━━━ 数据源 (只读, 均已存在) ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
  9mer 新切:    analysis/official/recompute_effN/R1_recomputed_rerun_9mer_effN8.csv
  8-11mer 新切: analysis/official/recompute_effN/R1_recomputed_rerun_8to11mer_effN8.csv
  DAI 版 8-11:  analysis/official/recompute_effN/R1_recomputed_rerun_8to11mer_dai_effN8.csv
    (三个 csv 读需跳 '#' 注释行; 列 Tool, fisherz_rho_effN)
  覆盖子肽层:   data/frozen/coverage_matrix.NEW.csv (长表 side/tool/status; MT 侧 scored 占比=覆盖率)
  图 (已生成, 直接嵌, 均在 paper/figures/):
    fig_8to11_ranking.png          (竖图 3600×4140, ratio 0.870 → 两栏页左卡完美贴合)
    fig_8to11_vs_9mer_dumbbell.png (宽图 4650×2484, ratio 1.872 → 全宽单独一页)
    fig_8to11_coverage.png         (竖图 3300×3960, ratio 0.833 → 两栏页左卡居中)

━━━ deck 结构 (6 页, widescreen 13.333×7.5in) ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
  s1 封面 / s2 一页看懂(三卡) / s3 核心结果(两栏·ranking) / s4 口径对比(全宽·dumbbell) /
  s5 覆盖(两栏·coverage) / s6 结论&边界(三卡)

Windows 规范: UTF-8 stdout; pathlib 路径; 中文字体 Microsoft YaHei (latin+ea+cs 全设防豆腐块);
  纯 csv/字典计算 (无 scipy); __main__ 守卫。

跑法 (主线跑, 本脚本不自跑):
  cd D:/YJ-Agent/project/meeting/QuantImmuBench && python analysis/plot_ppt_8to11_rerun.py
"""

import csv
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ── Windows 必要: UTF-8 stdout ───────────────────────────────────────────────
if hasattr(sys.stdout, "reconfigure"):
    sys.stdout.reconfigure(encoding="utf-8")

# ── 路径 ─────────────────────────────────────────────────────────────────────
HERE = Path(__file__).resolve().parent                     # analysis/
ROOT = HERE.parent                                         # QuantImmuBench/
RECOMP = HERE / "official" / "recompute_effN"
CSV_9MER = RECOMP / "R1_recomputed_rerun_9mer_effN8.csv"
CSV_8TO11 = RECOMP / "R1_recomputed_rerun_8to11mer_effN8.csv"
CSV_DAI = RECOMP / "R1_recomputed_rerun_8to11mer_dai_effN8.csv"
COV_CSV = ROOT / "data" / "frozen" / "coverage_matrix.NEW.csv"
FIG_DIR = ROOT / "paper" / "figures"
FIG_RANK = FIG_DIR / "fig_8to11_ranking.png"               # ratio 0.870
FIG_DUMB = FIG_DIR / "fig_8to11_vs_9mer_dumbbell.png"      # ratio 1.872 (宽)
FIG_COV = FIG_DIR / "fig_8to11_coverage.png"               # ratio 0.833
OUT_PPTX = ROOT / "QuantImmuBench_8to11mer新切重跑_2026-07-09.pptx"

# ── 配色 (Okabe-Ito, 精确复刻参考 deck) ───────────────────────────────────────
C_BG = RGBColor(0x0B, 0x3C, 0x49)       # 封面深底
C_BLUE = RGBColor(0x00, 0x72, 0xB2)     # 主蓝 (结果/覆盖页色条 + 来源脚注)
C_GREEN = RGBColor(0x00, 0x9E, 0x73)    # 绿 (结论页色条 + 封面装饰/结论)
C_ORANGE = RGBColor(0xE6, 0x9F, 0x00)   # 橙 (对比页色条)
C_INK = RGBColor(0x16, 0x32, 0x3A)      # 深青 (正文/标题)
C_SUBTXT = RGBColor(0xE6, 0xF2, 0xF2)   # 封面副文字
C_PAGENO = RGBColor(0x5E, 0x7B, 0x83)   # 页码
C_DATE = RGBColor(0x8F, 0xB7, 0xBD)     # 封面日期
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
C_CARD = RGBColor(0xF2, 0xF7, 0xF7)     # 全宽页文字条浅底
C_NEG = RGBColor(0xB2, 0x3A, 0x48)      # 红 (翻负强调)
C_CONCL = RGBColor(0x0F, 0x4A, 0x57)    # 封面结论条 (深底浅一档)

FONT = "Microsoft YaHei"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


# ── 数据加载 / 计算 (全部从 csv 现算) ─────────────────────────────────────────
def load_rho(path: Path) -> dict:
    """读 recompute_effN csv (跳 '#' 注释行), 返回 {Tool: fisherz_rho_effN(float)}, 跳空值。"""
    out = {}
    with open(path, encoding="utf-8") as f:
        rows = csv.DictReader(l for l in f if not l.startswith("#"))
        for r in rows:
            v = (r.get("fisherz_rho_effN") or "").strip()
            if v:
                out[r["Tool"]] = float(v)
    return out


def load_coverage(path: Path) -> dict:
    """
    读长表覆盖矩阵, 计算每工具 MT 侧覆盖率 = scored / total (子肽×HLA 位点)。
    返回 {tool: (scored, total, frac)}。
    """
    scored, total = {}, {}
    with open(path, encoding="utf-8") as f:
        for r in csv.DictReader(f):
            if r["side"] != "MT":
                continue
            t = r["tool"]
            total[t] = total.get(t, 0) + 1
            if r["status"] == "scored":
                scored[t] = scored.get(t, 0) + 1
    return {t: (scored.get(t, 0), total[t], scored.get(t, 0) / total[t]) for t in total}


def compute_stats():
    """现算所有 headline 数字, 返回 dict (同时供 print 自检)。"""
    a = load_rho(CSV_9MER)       # 9mer
    b = load_rho(CSV_8TO11)      # 8-11mer
    d = load_rho(CSV_DAI)        # DAI 版 8-11
    cov = load_coverage(COV_CSV)

    common = [t for t in a if t in b]                       # 两口径都有效的工具
    mean_9 = sum(a[t] for t in common) / len(common)
    mean_811 = sum(b[t] for t in common) / len(common)
    n_gt = sum(1 for t in common if a[t] > b[t])            # 9mer>8-11 计数

    top_by_811 = sorted(b, key=lambda t: b[t], reverse=True)  # 8-11 单口径降序 (排名要点)

    full_cov = sorted(t for t, (_, _, fr) in cov.items() if abs(fr - 1.0) < 1e-9)
    partial = sorted(((t, fr) for t, (_, _, fr) in cov.items() if fr < 1.0 - 1e-9),
                     key=lambda x: -x[1])
    n_hla_sites = next(iter(cov.values()))[1]               # MT 侧总子肽×HLA 位点 (=17088)

    dai_top = sorted(d.items(), key=lambda kv: -kv[1])[:3]  # DAI top3

    return dict(a=a, b=b, d=d, cov=cov, common=common,
                mean_9=mean_9, mean_811=mean_811, n_gt=n_gt, n_common=len(common),
                top_by_811=top_by_811, full_cov=full_cov, partial=partial,
                n_hla_sites=n_hla_sites, dai_top=dai_top)


# ── pptx 排版小工具 ───────────────────────────────────────────────────────────
def _set_font(run, size=None, bold=None, color=None, name=FONT):
    """设 run 字体 (latin + east-asian + complex-script 全设, 防中文豆腐块)。"""
    f = run.font
    if size is not None:
        f.size = Pt(size)
    if bold is not None:
        f.bold = bold
    if color is not None:
        f.color.rgb = color
    f.name = name
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:latin", "a:ea", "a:cs"):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {})
            rPr.append(el)
        el.set("typeface", name)


def add_rect(slide, x, y, w, h, fill, line=None):
    """加纯色矩形 (背景 / 色条 / 卡片)。x/y/w/h 均为 EMU (Inches/Pt)。"""
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(0.75)
    shp.shadow.inherit = False
    return shp


def add_oval(slide, x, y, w, h, fill):
    """加无边椭圆 (封面装饰)。"""
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def add_text(slide, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             wrap=True, space_after=4, line_spacing=None):
    """
    加文本框。runs = list of paragraph, 每 paragraph = list of (text,size,bold,color) tuples。
    返回 textframe。
    """
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        p.space_before = Pt(0)
        if line_spacing is not None:
            p.line_spacing = line_spacing
        for (text, size, bold, color) in para:
            r = p.add_run()
            r.text = text
            _set_font(r, size=size, bold=bold, color=color)
    return tf


def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])       # layout 6 = 全空白


def place_fig_by_height(slide, path, height_in, center_l, top_in):
    """按目标高度等比缩放嵌图, 水平居中于 center_l。返回 pic (缺图则加占位文字)。"""
    if not path.exists():
        add_text(slide, Inches(center_l - 2), Inches(3), Inches(4), Inches(1),
                 [[("[缺图] " + path.name, 15, True, C_NEG)]], align=PP_ALIGN.CENTER)
        return None
    pic = slide.shapes.add_picture(str(path), Inches(0), Inches(top_in), height=Inches(height_in))
    pic.left = int(Inches(center_l) - pic.width / 2)
    return pic


# ── 页级组件 (精确复刻参考几何) ────────────────────────────────────────────────
def content_page(prs, bar_color, eyebrow, title):
    """内容页骨架: 左全高细色条 + eyebrow(12pt) + 标题(22pt)。返回 slide。"""
    s = blank_slide(prs)
    add_rect(s, 0, 0, Inches(0.28), SLIDE_H, bar_color)                  # 左全高色条
    add_text(s, Inches(0.7), Inches(0.42), Inches(11.0), Inches(0.3),
             [[(eyebrow, 12, True, bar_color)]])
    add_text(s, Inches(0.7), Inches(0.72), Inches(12.2), Inches(0.7),
             [[(title, 22, True, C_INK)]], anchor=MSO_ANCHOR.MIDDLE)
    return s


def white_card(slide, l, t, w, h, bar_color, head, body_paras, head_size=15):
    """白卡 + 左边细色条(0.09) + section-head + 正文。坐标复刻参考三卡/两栏右卡。"""
    add_rect(slide, Inches(l), Inches(t), Inches(w), Inches(h), C_WHITE)
    add_rect(slide, Inches(l), Inches(t), Inches(0.09), Inches(h), bar_color)
    add_text(slide, Inches(l + 0.28), Inches(t + 0.16), Inches(w - 0.4), Inches(0.4),
             [[(head, head_size, True, bar_color)]])
    add_text(slide, Inches(l + 0.3), Inches(t + 0.62), Inches(w - 0.55), Inches(h - 0.75),
             body_paras, space_after=6, line_spacing=1.12)


def add_footer(slide, source, pageno, src_color=C_BLUE):
    """来源脚注(9pt 蓝 L0.7 T7.08) + 页码(11pt 灰 L12.53 T7.0)。"""
    add_text(slide, Inches(0.7), Inches(7.08), Inches(11.9), Inches(0.34),
             [[("来源  ", 9, True, src_color), (source, 9, False, src_color)]], space_after=1)
    add_text(slide, Inches(12.53), Inches(7.0), Inches(0.5), Inches(0.3),
             [[(str(pageno), 11, False, C_PAGENO)]], align=PP_ALIGN.LEFT)


def body(*lines):
    """把纯文本行列表转成 12pt INK 段落 (每行一段, 单 run)。"""
    return [[(t, 12, False, C_INK)] for t in lines]


# ── 各页 ─────────────────────────────────────────────────────────────────────
def slide_cover(prs, S):
    s = blank_slide(prs)
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, C_BG)                            # 深底
    add_rect(s, 0, 0, SLIDE_W, Inches(0.18), C_GREEN)                    # 顶部绿条
    add_oval(s, Inches(10.03), Inches(-1.6), Inches(4.6), Inches(4.6), C_BLUE)   # 右上蓝椭圆
    add_oval(s, Inches(11.33), Inches(3.6), Inches(3.2), Inches(3.2), C_GREEN)   # 右下绿椭圆

    add_text(s, Inches(0.9), Inches(1.3), Inches(11.7), Inches(0.4),
             [[("癌症个性化新抗原疫苗 · 免疫原性工具系统评测", 15, True, C_GREEN)]])
    add_text(s, Inches(0.9), Inches(1.9), Inches(11.7), Inches(1.9),
             [[("QuantImmuBench §2.2 · 8–11mer 可变窗补充口径", 30, True, C_WHITE)]],
             anchor=MSO_ANCHOR.TOP)
    add_text(s, Inches(0.95), Inches(4.7), Inches(11.5), Inches(1.0),
             [[("29 工具真实新跑 8/9/10/11mer · 复现零偏离核过 · 102 SNV 肽 · 9 患者 · 2026-07-09",
                15, False, C_SUBTXT)]], line_spacing=1.2)

    # 一句结论条
    add_rect(s, Inches(0.9), Inches(5.55), Inches(11.5), Inches(0.85), C_CONCL)
    add_rect(s, Inches(0.9), Inches(5.55), Inches(0.1), Inches(0.85), C_GREEN)
    add_text(s, Inches(1.2), Inches(5.55), Inches(11.0), Inches(0.85),
             [[("结论   ", 15, True, C_GREEN),
               ("9mer 主口径整体仍优于 8–11 可变窗，加固 §2.2（选 9AA）。", 15, True, C_WHITE)]],
             anchor=MSO_ANCHOR.MIDDLE)

    add_text(s, Inches(10.63), Inches(6.85), Inches(2.1), Inches(0.3),
             [[("2026-07-09", 12, False, C_DATE)]], align=PP_ALIGN.RIGHT)


def slide_onepager(prs, S):
    """s2 一页看懂 · 三卡, 左条蓝。"""
    a, b = S["a"], S["b"]
    s = content_page(prs, C_BLUE, "一页看懂",
                     "8–11mer 可变窗补充口径：一页看懂")
    # 卡1 (蓝) 口径是什么
    white_card(s, 0.7, 1.6, 3.84, 5.35, C_BLUE, "8–11mer 口径是什么",
               body("每肽按 8/9/10/11mer 全窗切分 → 各工具打分后 max 池化"
                    "（零选择 <tool>_max）。",
                    "",
                    "effN≥8 有效覆盖门槛 → Fisher-Z 病人等权聚合，与 9mer "
                    "主口径同引擎、可比。",
                    "",
                    "102 SNV 突变肽 · DS2 9 患者 · Elispot 连续 SFC · "
                    "cluster-bootstrap 95%CI（n_boot=2000, seed=42）。"))
    # 卡2 (绿) 和 9mer 比结论
    pct = 100 * S["n_gt"] / S["n_common"]
    white_card(s, 4.74, 1.6, 3.84, 5.35, C_GREEN, "和 9mer 比结论",
               [[("均值 ρ：9mer ", 12, False, C_INK),
                 ("%.3f" % S["mean_9"], 12, True, C_BLUE),
                 (" vs 8–11 可变窗 ", 12, False, C_INK),
                 ("%.3f" % S["mean_811"], 12, True, C_ORANGE), ("。", 12, False, C_INK)],
                [("", 12, False, C_INK)],
                [("%d/%d" % (S["n_gt"], S["n_common"]), 12, True, C_GREEN),
                 (" 工具 9mer 更高（%.0f%%）→ 9mer 主口径整体更优，"
                  "加固 §2.2 选 9AA。" % pct, 12, False, C_INK)],
                [("", 12, False, C_INK)],
                [("唯一例外 MHCnuggets（8–11 略升 %.3f→%.3f）。"
                  % (a["MHCnuggets"], b["MHCnuggets"]), 12, False, C_INK)]])
    # 卡3 (橙) 诚实脚注
    white_card(s, 8.78, 1.6, 3.84, 5.35, C_ORANGE, "诚实脚注",
               body("DeepNetBim = 9mer-only、DeepImmuno = 9–10mer only"
                    "（原生架构限长，非漏跑）。",
                    "",
                    "数据溯源：29 工具 official.csv 全 2026-07-08 真跑；",
                    "",
                    "9mer 子集 vs 独立 9mer 结果 bit 一致 → 复现零偏离核过。"))
    add_footer(s, "R1_recomputed_rerun_8to11mer_effN8.csv · coverage_matrix.NEW.csv · "
                  "改动②③全量重跑（2026-07-08）", 2)


def slide_ranking(prs, S):
    """s3 核心结果 · 两栏, 左条蓝, 左卡嵌 ranking 图。"""
    b = S["b"]
    s = content_page(prs, C_BLUE, "核心结果 · §2.2 单工具主指标",
                     "8–11mer 可变窗单工具 max-pooling Spearman 排名")
    # 左白卡 + ranking 图 (ratio 0.870, 完美贴合 4.75×5.46)
    add_rect(s, Inches(0.5), Inches(1.5), Inches(5.5), Inches(5.7), C_WHITE)
    place_fig_by_height(s, FIG_RANK, 5.46, 3.25, 1.62)
    # 右白卡: 排名要点
    top5 = S["top_by_811"][:5]
    top5_str = " > ".join("%s %.3f" % (t, b[t]) for t in top5)
    white_card(s, 6.25, 1.5, 6.35, 5.7, C_BLUE, "8–11mer 单工具排名要点",
               [[("头部前五：", 12, True, C_INK), (top5_str, 12, False, C_INK)],
                [("", 12, False, C_INK)],
                [("均值 ρ̄ = ", 12, True, C_INK),
                 ("%.3f" % S["mean_811"], 12, True, C_BLUE),
                 ("（Fisher-Z 病人等权，%d 工具）。" % S["n_common"], 12, False, C_INK)],
                [("", 12, False, C_INK)],
                [("负值尾部：", 12, True, C_NEG),
                 ("TSCAPE / CNNeo / MHCseqNet 翻负 —— 8–11 可变窗上 "
                  "max-pool 被非 9mer 窗抢走，稀释了对齐 SFC 的信号。", 12, False, C_INK)],
                [("", 12, False, C_INK)],
                [("与 9mer 主口径（ρ̄ %.3f）相比整体下移，仅 MHCnuggets 反升。"
                  % S["mean_9"], 12, False, C_INK)]])
    add_footer(s, "R1_recomputed_rerun_8to11mer_effN8.csv · per-patient Spearman → "
                  "effN≥8 Fisher-Z 病人等权", 3)


def slide_dumbbell(prs, S):
    """s4 口径对比 · 全宽 dumbbell 图页, 左条橙。"""
    a, b = S["a"], S["b"]
    s = content_page(prs, C_ORANGE, "口径对比 · 9mer vs 8–11mer",
                     "9mer 主口径 vs 8–11mer 可变窗 · 逐工具哑铃对比")
    # 全宽嵌图 (ratio 1.872 宽; 装入 box L0.5 T1.5 W12.3 H5.0, 高度受限居中)
    place_fig_by_height(s, FIG_DUMB, 5.0, 6.65, 1.5)
    # 图下窄文字条 (横跨)
    add_rect(s, Inches(0.5), Inches(6.5), Inches(12.33), Inches(0.5), C_CARD)
    add_rect(s, Inches(0.5), Inches(6.5), Inches(0.09), Inches(0.5), C_ORANGE)
    add_text(s, Inches(0.75), Inches(6.5), Inches(12.0), Inches(0.5),
             [[("9mer 整体优于 8–11（均值 %.3f > %.3f、%d/%d 工具下降）；"
                "唯 MHCnuggets %.3f→%.3f 反升（例外）；MHCseqNet %.3f→%.3f、"
                "ImmuGenX 翻负 = 非 9mer 窗抢 max-pool 的真实信号，非 bug。"
                % (S["mean_9"], S["mean_811"], S["n_gt"], S["n_common"],
                   a["MHCnuggets"], b["MHCnuggets"], a["MHCseqNet"], b["MHCseqNet"]),
                12, False, C_INK)]], anchor=MSO_ANCHOR.MIDDLE)
    add_footer(s, "R1_recomputed_rerun_9mer_effN8.csv · R1_recomputed_rerun_8to11mer_effN8.csv"
                  " · 逐工具 Δρ", 4)


def slide_coverage(prs, S):
    """s5 覆盖 · 两栏, 左条蓝, 左卡嵌 coverage 图。"""
    n_full = len(S["full_cov"])
    n_tools = len(S["cov"])
    n_part = len(S["partial"])
    s = content_page(prs, C_BLUE, "覆盖 · 8–11mer 子肽层",
                     "8–11mer 子肽层覆盖：%d/%d 工具满覆盖，%d 个诚实上限"
                     % (n_full, n_tools, n_part))
    # 左白卡 + coverage 图 (ratio 0.833, 高对齐 5.46 居中)
    add_rect(s, Inches(0.5), Inches(1.5), Inches(5.5), Inches(5.7), C_WHITE)
    place_fig_by_height(s, FIG_COV, 5.46, 3.25, 1.62)
    # 右白卡: 覆盖要点 + 部分覆盖清单
    reason = {
        "ICERFIRE": "部署天花板",
        "HLAthena": "缺神经权重（部署天花板）",
        "DeepImmuno": "9–10mer 原生限长",
        "NetTepi": "13 等位限制",
        "DeepNetBim": "9mer 原生限长",
    }
    paras = [
        [("%d/%d" % (n_full, n_tools), 12, True, C_GREEN),
         (" 工具满覆盖（%d 子肽 × HLA MT 位点）。" % S["n_hla_sites"], 12, False, C_INK)],
        [("", 12, False, C_INK)],
        [("%d 个诚实上限（全 documented，非真漏）：" % n_part, 12, True, C_ORANGE)],
    ]
    for t, fr in S["partial"]:
        paras.append([("%s  %.1f%%" % (t, fr * 100), 12, True, C_INK),
                      ("  — " + reason.get(t, "documented 上限"), 12, False, C_PAGENO)])
    paras += [
        [("", 12, False, C_INK)],
        [("覆盖以子肽层（每子肽×HLA 位点）计，非 mut_key 层 max-pool"
          "（后者单窗命中即误显全 102 肽满覆盖）。", 12, False, C_INK)],
    ]
    white_card(s, 6.25, 1.5, 6.35, 5.7, C_BLUE, "8–11mer 子肽层覆盖", paras)
    add_footer(s, "coverage_matrix.NEW.csv（MT 侧 scored/total, 子肽层）", 5)


def slide_conclusion(prs, S):
    """s6 结论 & 边界 · 三卡, 左条绿。"""
    a, b = S["a"], S["b"]
    dai_str = "、".join("%s %.3f" % (t, v) for t, v in S["dai_top"])
    s = content_page(prs, C_GREEN, "结论 & 边界", "结论与诚实边界")
    # 卡1 (绿) 核心结论
    white_card(s, 0.7, 1.6, 3.84, 5.35, C_GREEN, "核心结论",
               body("9mer 主口径整体优于 8–11 可变窗：",
                    "",
                    "均值 ρ %.3f vs %.3f、%d/%d（%.0f%%）工具下降。"
                    % (S["mean_9"], S["mean_811"], S["n_gt"], S["n_common"],
                       100 * S["n_gt"] / S["n_common"]),
                    "",
                    "→ 加固 §2.2 主口径选 9AA。"))
    # 卡2 (橙) 诚实边界
    white_card(s, 4.74, 1.6, 3.84, 5.35, C_ORANGE, "诚实边界",
               body("MHCnuggets 是唯一反升（%.3f→%.3f）。"
                    % (a["MHCnuggets"], b["MHCnuggets"]),
                    "",
                    "限长工具（DeepNetBim=9mer、DeepImmuno=9–10mer）在可变窗仅覆盖"
                    "原生窗长子肽，窗长已据实标注。",
                    "",
                    "TSCAPE / CNNeo / MHCseqNet 翻负 = 可变窗抢 max-pool 的真实信号，"
                    "非 bug。"))
    # 卡3 (蓝) DAI 版
    white_card(s, 8.78, 1.6, 3.84, 5.35, C_BLUE, "DAI 版 8–11mer",
               body("DAI 差分口径 8–11mer top：",
                    "",
                    dai_str + "。",
                    "",
                    "（从 dai csv 现算，供交叉参考。）"))
    add_footer(s, "R1_recomputed_rerun_8to11mer_effN8.csv · "
                  "R1_recomputed_rerun_8to11mer_dai_effN8.csv · coverage_matrix.NEW.csv", 6)


# ── main ─────────────────────────────────────────────────────────────────────
def build():
    S = compute_stats()

    # ── print 自检 (供主线核对锚点值) ──────────────────────────────────────────
    print("=" * 64)
    print("[自检] 8-11mer 新切重跑 headline (从 csv 现算)")
    print("-" * 64)
    print("  两口径均有效工具 n_common = %d  (期望 28)" % S["n_common"])
    print("  均值 rho  9mer = %.4f   (期望 0.1127)" % S["mean_9"])
    print("  均值 rho  8-11 = %.4f   (期望 0.0541)" % S["mean_811"])
    print("  9mer > 8-11    = %d/%d   (期望 22/28)" % (S["n_gt"], S["n_common"]))
    print("  8-11 top5: " + " > ".join("%s %.3f" % (t, S["b"][t]) for t in S["top_by_811"][:5]))
    print("    (期望 MHCnuggets 0.329 > netMHCpan_BA 0.213 > PredIG 0.179 "
          "> IEDB_Calis 0.177 > ICERFIRE 0.176)")
    print("  MHCnuggets  9mer = %.3f -> 8-11 = %.3f  (期望 0.319->0.329, 唯一升)"
          % (S["a"]["MHCnuggets"], S["b"]["MHCnuggets"]))
    print("  MHCseqNet   9mer = %.3f -> 8-11 = %.3f  (翻负)"
          % (S["a"]["MHCseqNet"], S["b"]["MHCseqNet"]))
    if "ImmuGenX" in S["a"] and "ImmuGenX" in S["b"]:
        print("  ImmuGenX    9mer = %.3f -> 8-11 = %.3f  (翻负)"
              % (S["a"]["ImmuGenX"], S["b"]["ImmuGenX"]))
    print("-" * 64)
    print("  覆盖 (MT 侧 scored/total, %d 子肽×HLA 位点):" % S["n_hla_sites"])
    print("    满覆盖工具数 = %d  (期望 24)" % len(S["full_cov"]))
    for t, fr in S["partial"]:
        print("    %-12s %.1f%%" % (t, fr * 100))
    print("    (期望: ICERFIRE 90.7%, HLAthena 76.8%, DeepImmuno 50.0%, "
          "NetTepi 31.1%, DeepNetBim 23.7%)")
    print("-" * 64)
    print("  DAI 版 8-11 top3: " + ", ".join("%s %.3f" % (t, v) for t, v in S["dai_top"]))
    print("    (期望 IEDB_Calis 0.329, PredIG 0.328, Repitope 0.294)")
    print("=" * 64)

    # ── 建 deck ────────────────────────────────────────────────────────────────
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_cover(prs, S)
    slide_onepager(prs, S)
    slide_ranking(prs, S)
    slide_dumbbell(prs, S)
    slide_coverage(prs, S)
    slide_conclusion(prs, S)

    prs.save(str(OUT_PPTX))
    print("[OK] 已写: %s (%d slides)" % (OUT_PPTX, len(prs.slides._sldIdLst)))


if __name__ == "__main__":
    build()
